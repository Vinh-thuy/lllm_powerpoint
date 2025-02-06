import os
import collections.abc
import pptx

class TemplateProcessor:
    def __init__(self, template_path):
        try:
            self.presentation = pptx.Presentation(template_path)
        except Exception as e:
            print(f"Erreur lors du chargement du template : {e}")
            # Créer une nouvelle présentation si le template ne peut pas être chargé
            self.presentation = pptx.Presentation()

    def update_slide(self, slide_id, updates):
        try:
            # Vérifier si le slide existe
            if slide_id < 0 or slide_id >= len(self.presentation.slides):
                print(f"Slide {slide_id} n'existe pas")
                return

            slide = self.presentation.slides[slide_id]
            
            for update in updates:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text_frame = shape.text_frame
                        # Mise à jour du texte si spécifié
                        if 'text' in update:
                            text_frame.text = update['text']
        except Exception as e:
            print(f"Erreur lors de la mise à jour du slide {slide_id}: {e}")

    def save(self, output_path):
        try:
            # Créer le dossier de sortie s'il n'existe pas
            os.makedirs(os.path.dirname(output_path), exist_ok=True)
            self.presentation.save(output_path)
            print(f"Présentation sauvegardée : {output_path}")
        except Exception as e:
            print(f"Erreur lors de la sauvegarde : {e}")

    def create_new_slide(self, layout_index=0):
        """
        Créer un nouveau slide avec un layout spécifique
        :param layout_index: Index du layout (0 par défaut)
        :return: Le nouveau slide créé
        """
        try:
            slide_layout = self.presentation.slide_layouts[layout_index]
            return self.presentation.slides.add_slide(slide_layout)
        except Exception as e:
            print(f"Erreur lors de la création du slide : {e}")
            return None
