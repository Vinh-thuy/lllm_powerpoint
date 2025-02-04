import os
from dotenv import load_dotenv
from openai import OpenAI
import yaml
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import base64
import imghdr
from PIL import Image

# Charger les variables d'environnement du fichier .env
load_dotenv()

def load_config():
    with open('config.yaml', 'r') as f:
        return yaml.safe_load(f)

def encode_image_to_base64(image_path):
    if not os.path.exists(image_path):
        raise FileNotFoundError(f"L'image {image_path} n'existe pas. Veuillez vérifier le chemin.")
    
    # Vérifier la taille et le type de fichier
    file_size = os.path.getsize(image_path)
    file_type = imghdr.what(image_path)
    
    print(f"Fichier image : {image_path}")
    print(f"Taille du fichier : {file_size} octets")
    print(f"Type d'image : {file_type}")
    
    if file_size > 20 * 1024 * 1024:  # Limite à 20 Mo
        raise ValueError("La taille de l'image est trop grande (> 20 Mo)")
    
    if file_type not in ['jpeg', 'png', 'gif', 'bmp']:
        raise ValueError(f"Type de fichier non supporté : {file_type}")
    
    with open(image_path, "rb") as image_file:
        return base64.b64encode(image_file.read()).decode('utf-8')

def analyze_image_with_vision(client, image_path):
    try:
        base64_image = encode_image_to_base64(image_path)
        
        response = client.chat.completions.create(
            model="gpt-4o",  # Utilisation du modèle GPT-4o qui supporte la vision
            messages=[
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": "Analyse cette roadmap et donne-moi les informations suivantes en format JSON :\n"
                                   "1. Les mois présents\n"
                                   "2. Les tâches avec leurs périodes\n"
                                   "3. Les couleurs utilisées (en RGB)\n"
                                   "4. La légende (Done, Coming, Not planned)\n"
                        },
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:image/jpeg;base64,{base64_image}"
                            }
                        }
                    ]
                }
            ],
            max_tokens=1000
        )
        return response.choices[0].message.content
    except Exception as e:
        print(f"Erreur détaillée lors de l'analyse de l'image : {type(e).__name__} - {str(e)}")
        return None

def create_roadmap_slide(prs, roadmap_data=None):
    # Utiliser le layout blank
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Titre "ROADMAP"
    title = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(0.5))
    title.text_frame.text = "ROADMAP"
    title.text_frame.paragraphs[0].font.size = Pt(24)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 100, 50)
    
    # Grille des mois
    months_box = slide.shapes.add_table(2, 12, Inches(1), Inches(1.5), Inches(8), Inches(0.5)).table
    months = ["JAN", "FEB", "MAR", "APR", "MAY", "JUN", "JUL", "AUG", "SEP", "OCT", "NOV", "DEC"]
    
    for i, month in enumerate(months):
        cell = months_box.cell(0, i)
        cell.text = month
        cell.text_frame.paragraphs[0].font.size = Pt(8)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Ajouter les tâches
    task_y = 2.5
    task1 = slide.shapes.add_shape(1, Inches(1), Inches(task_y), Inches(3), Inches(0.4))
    task1.fill.solid()
    task1.fill.fore_color.rgb = RGBColor(0, 100, 50)
    task1.text_frame.text = "Task1"
    task1.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    task2 = slide.shapes.add_shape(1, Inches(2.5), Inches(task_y + 0.5), Inches(4), Inches(0.4))
    task2.fill.solid()
    task2.fill.fore_color.rgb = RGBColor(200, 255, 200)
    task2.text_frame.text = "Task2"
    
    # Ajouter le logo BNP
    logo_left = Inches(1)
    logo_bottom = Inches(6.5)
    logo = slide.shapes.add_textbox(logo_left, logo_bottom, Inches(4), Inches(0.5))
    logo.text_frame.text = "BNP PARIBAS"
    logo.text_frame.paragraphs[0].font.bold = True
    
    # Ajouter le texte "The bank for a changing world"
    tagline = slide.shapes.add_textbox(Inches(4), logo_bottom, Inches(4), Inches(0.5))
    tagline.text_frame.text = "The bank for a changing world"
    
    return slide

def main():
    config = load_config()
    
    # Récupérer la clé API depuis les variables d'environnement
    api_key = os.getenv('OPENAI_API_KEY')
    if not api_key:
        print("Erreur : La clé API OpenAI n'est pas définie dans le fichier .env")
        return
    
    client = OpenAI(api_key=api_key)
    
    # Créer une nouvelle présentation
    prs = Presentation()
    
    # Définir la taille des slides en 16:9
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # Chemin de l'image fixe
    image_path = "/Users/vinh/Documents/LLM_PPT/img_src/roadmap.jpg"
    
    # Vérification détaillée de l'image
    if not os.path.exists(image_path):
        print(f"ERREUR CRITIQUE : Le fichier {image_path} n'existe pas !")
        return
    
    print(f"Chemin de l'image : {image_path}")
    print(f"Taille du fichier : {os.path.getsize(image_path)} octets")
    
    try:
        with Image.open(image_path) as img:
            print(f"Dimensions de l'image : {img.width}x{img.height}")
            print(f"Format de l'image : {img.format}")
            print(f"Mode de l'image : {img.mode}")
    except Exception as e:
        print(f"Erreur lors de l'ouverture de l'image : {e}")
    
    # Analyse de l'image
    roadmap_data = analyze_image_with_vision(client, image_path)
    
    # Créer la slide de roadmap
    create_roadmap_slide(prs, roadmap_data)
    
    # Sauvegarder la présentation
    output_path = os.path.join(config['output_dir'], 'roadmap.pptx')
    os.makedirs(config['output_dir'], exist_ok=True)
    prs.save(output_path)
    print(f"Présentation générée : {output_path}")

if __name__ == "__main__":
    main()
