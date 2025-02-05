import os
from dotenv import load_dotenv
import ollama
import yaml
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import MSO_SHAPE
import base64
import imghdr
from PIL import Image
import re
from datetime import datetime
import json
import traceback
import sys

# Charger les variables d'environnement du fichier .env
load_dotenv()

def load_config():
    with open('config.yaml', 'r') as f:
        return yaml.safe_load(f)

def encode_image_to_base64(image_path):
    if not os.path.exists(image_path):
        raise FileNotFoundError(f"L'image {image_path} n'existe pas. Veuillez vérifier le chemin.")
    
    # Vérifier la taille et le type de fichier
    file_size = os.path.getsize(image_path)
    file_type = imghdr.what(image_path)
    
    print(f"Fichier image : {image_path}")
    print(f"Taille du fichier : {file_size} octets")
    print(f"Type d'image : {file_type}")
    
    if file_size > 20 * 1024 * 1024:  # Limite à 20 Mo
        raise ValueError("La taille de l'image est trop grande (> 20 Mo)")
    
    if file_type not in ['jpeg', 'png', 'gif', 'bmp']:
        raise ValueError(f"Type de fichier non supporté : {file_type}")
    
    with open(image_path, "rb") as image_file:
        return base64.b64encode(image_file.read()).decode('utf-8')

def analyze_image_with_vision(client, image_path):
    try:
        base64_image = encode_image_to_base64(image_path)
        
        response = client.chat(
            model=config['llm_model'],  
            messages=[
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": "Analyse cette roadmap et donne-moi les informations suivantes en format JSON :\n"
                                   "1. Les mois présents\n"
                                   "2. Les tâches avec leurs périodes\n"
                                   "3. Les couleurs utilisées (en RGB)\n"
                                   "4. La légende (Done, Coming, Not planned)\n"
                        },
                        {
                            "type": "image",
                            "image": base64_image
                        }
                    ]
                }
            ],
            max_tokens=1000
        )
        return response.choices[0].message.content
    except Exception as e:
        print(f"Erreur détaillée lors de l'analyse de l'image : {type(e).__name__} - {str(e)}")
        return None

def parse_date_to_month_position(client, date_str, config=None):
    """
    Convertit une date en position de mois sur la timeline
    
    Args:
        client: Client Ollama
        date_str (str): Chaîne représentant la date
        config (dict, optional): Configuration du projet
    
    Returns:
        tuple: (mois, position_dans_le_mois)
    """
    # Dictionnaire de mapping des mois
    mois_mapping = {
        'janvier': 0, 'jan': 0, 
        'février': 1, 'fevrier': 1, 'fev': 1,
        'mars': 2, 'mar': 2,
        'avril': 3, 'avr': 3,
        'mai': 4,
        'juin': 5, 'jun': 5,
        'juillet': 6, 'jul': 6,
        'août': 7, 'aout': 7, 'aug': 7,
        'septembre': 8, 'sep': 8,
        'octobre': 9, 'oct': 9,
        'novembre': 10, 'nov': 10,
        'décembre': 11, 'decembre': 11, 'dec': 11
    }
    
    # Nettoyer et normaliser la date
    date_str = date_str.lower().strip()
    
    # Extraire le jour et le mois
    jour_match = re.search(r'(\d+)', date_str)
    mois_match = re.search(r'(janvier|jan|février|fevrier|fev|mars|avril|avr|mai|juin|jul|août|aout|septembre|sep|octobre|oct|novembre|nov|décembre|decembre|dec)', date_str, re.IGNORECASE)
    
    # Valeurs par défaut
    jour = 1
    mois_index = 0
    
    # Extraire le jour si présent
    if jour_match:
        try:
            jour = int(jour_match.group(1))
        except ValueError:
            jour = 1
    
    # Extraire le mois
    if mois_match:
        mois = mois_match.group(1).lower()
        for key, value in mois_mapping.items():
            if mois in key:
                mois_index = value
                break
    
    # Calculer la position dans le mois (début, milieu, fin)
    if jour <= 10:
        position_mois = 0.0  # début du mois
    elif jour <= 20:
        position_mois = 0.5  # milieu du mois
    else:
        position_mois = 1.0  # fin du mois
    
    print(f"Analyse de la date : {date_str}")
    print(f"  Jour extrait : {jour}")
    print(f"  Mois extrait : {mois}")
    print(f"  Index du mois : {mois_index}")
    print(f"  Position dans le mois : {position_mois}")
    
    return (mois_index, position_mois)

def parse_prompt_with_llm(client, prompt, config):
    """Analyse le prompt pour extraire les informations de la tâche"""
    # Liste des modèles à essayer
    models_to_try = [
        config.get('llm_model', 'mervinpraison/llama3.2-3B-instruct-test-2:8b'),
    ]
    
    for model in models_to_try:
        try:
            print(f"Tentative avec le modèle : {model}")
            
            response = client.chat(
                model=model,
                messages=[
                    {
                        'role': 'system', 
                        'content': """Tu es un assistant spécialisé dans l'analyse de prompts de projet. 
                        Pour chaque prompt, réponds avec les informations suivantes :
                        - Nom du projet
                        - Date de début
                        - Date de fin
                        - Couleur (optionnelle)
                        
                        Exemple de format de réponse :
                        Projet: [Nom du projet]
                        Début: [Date de début]
                        Fin: [Date de fin]
                        Couleur: [Couleur (optionnelle)]"""
                    },
                    {
                        'role': 'user', 
                        'content': prompt
                    }
                ]
            )
            
            result = response['message']['content'].strip()
            print("Réponse du modèle :", result)  # Ajout de debug
            
            # Extraction des informations de base
            name_match = re.search(r"Projet:\s*([^\n]+)", result, re.IGNORECASE)
            start_date_match = re.search(r"Début:\s*([^\n]+)", result, re.IGNORECASE)
            end_date_match = re.search(r"Fin:\s*([^\n]+)", result, re.IGNORECASE)
            color_match = re.search(r"Couleur:\s*([^\n]+)", result, re.IGNORECASE)
            
            if not (name_match and start_date_match and end_date_match):
                # Fallback à l'extraction depuis le prompt original
                name_match = re.search(r"projet '([^']+)'", prompt)
                start_date_match = re.search(r"du (\d+\s*\w+)", prompt)
                end_date_match = re.search(r"au (\d+\s*\w+)", prompt)
                color_match = re.search(r"couleur\s*:\s*(\w+)", prompt, re.IGNORECASE)
            
            if not (name_match and start_date_match and end_date_match):
                print(f"Format de prompt invalide pour le modèle {model}")
                continue
            
            # Extraire les informations
            task_name = name_match.group(1).strip()
            start_date = start_date_match.group(1).strip()
            end_date = end_date_match.group(1).strip()
            
            # Couleur par défaut
            color = color_match.group(1).lower().strip() if color_match else "bleu"
            
            # Convertir les dates en positions de mois
            start_month = parse_date_to_month_position(client, start_date, config)
            end_month = parse_date_to_month_position(client, end_date, config)
            
            # Mapping des couleurs
            color_map = {
                'rouge': [255, 0, 0],
                'bleu': [0, 0, 255],
                'vert': [0, 255, 0],
                'jaune': [255, 255, 0],
                'orange': [255, 165, 0],
                'violet': [128, 0, 128],
                'rose': [255, 192, 203],
                'marron': [165, 42, 42]
            }
            
            return {
                'task_name': task_name,
                'start_month': start_month,
                'end_month': end_month,
                'color_rgb': color_map.get(color, [0, 0, 255])  # Bleu par défaut
            }
        
        except Exception as e:
            print(f"Erreur avec le modèle {model}: {e}")
            continue
    
    print("Aucun modèle n'a pu traiter le prompt")
    return None

def create_task_on_roadmap(prs, task_info):
    """
    Crée une tâche sur la roadmap PowerPoint
    
    Args:
        prs (Presentation): Présentation PowerPoint
        task_info (dict): Informations de la tâche
    """
    # Récupérer le slide de roadmap (premier slide)
    slide = prs.slides[0]
    
    # Dimensions de la slide
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    
    # Dimensions effectives de la grille
    grid_margin_x = Inches(0.5)
    grid_width = slide_width - (2 * grid_margin_x)
    
    # Calculer les positions de début et de fin
    start_month, start_pos = task_info['start_month']
    end_month, end_pos = task_info['end_month']
    
    # Ajuster la position précise dans le mois
    month_width = grid_width / 12
    
    start_x = grid_margin_x + (start_month * month_width) + (start_pos * month_width)
    end_x = grid_margin_x + ((end_month + 1) * month_width) - ((1 - end_pos) * month_width)
    
    # Largeur de la tâche
    task_width = max(end_x - start_x, Inches(0.5))  # Largeur minimale
    
    # Calculer la position Y
    existing_shapes = [shape for shape in slide.shapes if shape.has_text_frame]
    y_position = Inches(2.5 + (len(existing_shapes) * 0.6))
    
    # Créer la forme de la tâche
    task_shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 
        start_x, 
        y_position, 
        task_width, 
        Inches(0.4)
    )
    
    # Couleur de la tâche
    task_shape.fill.solid()
    task_shape.fill.fore_color.rgb = RGBColor(*task_info['color_rgb'])
    
    # Ajouter le texte de la tâche
    text_frame = task_shape.text_frame
    text_frame.text = task_info['task_name']
    text_frame.paragraphs[0].font.size = Pt(12)
    text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # Noir
    
    print(f"Tâche '{task_info['task_name']}' créée avec succès !")
    print(f"  Mois de début : {start_month}, Position : {start_pos}")
    print(f"  Mois de fin : {end_month}, Position : {end_pos}")
    print(f"  Position X de début : {start_x / 914400:.2f} inches")
    print(f"  Position X de fin : {end_x / 914400:.2f} inches")
    print(f"  Largeur de la tâche : {task_width / 914400:.2f} inches")

def create_roadmap_slide(prs, task_info=None):
    """Crée ou met à jour un slide de roadmap"""
    print("\n=== Création/Mise à jour du slide de roadmap ===")
    
    # Utiliser le premier slide ou en créer un nouveau
    if len(prs.slides) > 0:
        slide = prs.slides[0]
        print(f"Utilisation du slide existant (index 0)")
    else:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Layout vide
        print(f"Création d'un nouveau slide")
    
    # Vérifier si la grille des mois existe déjà
    months_grid_exists = any(shape.has_table for shape in slide.shapes)
    
    # Titre "ROADMAP"
    if not any(shape.has_text_frame and shape.text_frame.text == "ROADMAP" for shape in slide.shapes):
        title = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(0.5))
        title.text_frame.text = "ROADMAP"
        title.text_frame.paragraphs[0].font.size = Pt(24)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)    
    
    if not months_grid_exists:
        print("Ajout de la grille des mois")
        
        # Calculer les dimensions de la slide
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        
        print(f"Dimensions de la slide : {slide_width / 914400:.2f} x {slide_height / 914400:.2f} inches")
        
        # Marges internes
        margin_left = Inches(0.5)
        margin_right = Inches(0.5)
        
        # Calculer la largeur effective
        effective_width = slide_width - (margin_left + margin_right)
        
        # Créer la grille des mois
        months_box = slide.shapes.add_table(
            2,  # 2 rangées
            12,  # 12 colonnes (mois)
            margin_left,  # Position X de départ
            Inches(1.5),  # Position Y
            effective_width,  # Largeur totale
            Inches(0.5)  # Hauteur
        ).table
        
        months = ["JAN", "FEB", "MAR", "APR", "MAY", "JUN", "JUL", "AUG", "SEP", "OCT", "NOV", "DEC"]
        
        for i, month in enumerate(months):
            cell = months_box.cell(0, i)
            cell.text = month
            cell.text_frame.paragraphs[0].font.size = Pt(8)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    else:
        print("Grille des mois déjà existante")
    
    # Ajouter la nouvelle tâche si spécifiée
    if task_info:
        create_task_on_roadmap(prs, task_info)
    
    return slide

def main():
    # Créer les dossiers de sortie
    templates_dir = "templates"
    output_dir = "generated"
    
    # Créer les dossiers s'ils n'existent pas
    os.makedirs(templates_dir, exist_ok=True)
    os.makedirs(output_dir, exist_ok=True)
    
    # Chemins complets
    template_path = os.path.join(templates_dir, "roadmap_template.pptx")
    output_path = os.path.join(output_dir, "roadmap.pptx")
    
    # Supprimer le fichier existant dans le répertoire de génération
    if os.path.exists(output_path):
        try:
            os.remove(output_path)
            print(f"Fichier existant supprimé : {output_path}")
        except Exception as e:
            print(f"Erreur lors de la suppression du fichier : {e}")
    
    # Charger ou créer la présentation
    if os.path.exists(template_path):
        prs = Presentation(template_path)
        print(f"Chargement du template depuis {template_path}")
    else:
        prs = Presentation()
        print("Création d'une nouvelle présentation")
        
        # Sauvegarde du template vide
        prs.save(template_path)
        print(f"Template vide sauvegardé dans {template_path}")
    
    # Configuration
    config = load_config()
    client = ollama.Client(host=config.get('ollama_host', 'http://localhost:11434'))
    
    # Lire les prompts
    prompts = []
    
    # Vérifier si des arguments en ligne de commande sont passés
    if len(sys.argv) > 1:
        # Concaténer tous les arguments en ligne de commande
        prompts = [' '.join(sys.argv[1:])]
    
    # Si pas d'arguments en ligne de commande, essayer de lire le fichier
    if not prompts:
        try:
            with open('prompt.txt', 'r') as f:
                prompts = [line.strip() for line in f.readlines() if line.strip()]
        except FileNotFoundError:
            # Si le fichier n'existe pas, demander une saisie interactive
            prompts = [input("Entrez votre demande (ex: ajoute le projet TOTO se déroule du 1er février au 13 juin (couleur : orange)) : ")]
    
    # Traiter chaque prompt successivement
    for prompt in prompts:
        print(f"\n--- Traitement du prompt : {prompt} ---")
        
        try:
            task_info = parse_prompt_with_llm(client, prompt, config)
            
            if task_info:
                # Créer ou mettre à jour le slide de roadmap
                create_roadmap_slide(prs, task_info)
                
                # Sauvegarder la présentation après chaque prompt
                prs.save(output_path)
                print(f"Présentation mise à jour : {output_path}")
            else:
                print("Impossible de parser le prompt")
        
        except Exception as e:
            print(f"Erreur lors du traitement du prompt '{prompt}' : {e}")
            traceback.print_exc()
    
    print("\nTraitement de tous les prompts terminé.")

if __name__ == "__main__":
    main()