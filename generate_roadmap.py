import os
from dotenv import load_dotenv
import ollama
import yaml
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import MSO_SHAPE
import base64
import imghdr
from PIL import Image
import re
from datetime import datetime
import json
import traceback
import sys

# Charger les variables d'environnement du fichier .env
load_dotenv()

# Initialiser la base de données
from task_database import TaskDatabase
task_db = TaskDatabase()

def convert_color_to_rgb(color):
    """
    Convertit un nom de couleur en code RGB
    
    Args:
        color (str): Nom de la couleur
    
    Returns:
        list: Code RGB de la couleur
    """
    color_map = {
        'rouge': [255, 0, 0],
        'bleu': [0, 0, 255],
        'vert': [0, 255, 0],
        'jaune': [255, 255, 0],
        'orange': [255, 165, 0],
        'violet': [128, 0, 128],
        'rose': [255, 192, 203],
        'marron': [165, 42, 42],
        'gris': [128, 128, 128],
        'noir': [0, 0, 0],
        'blanc': [255, 255, 255]
    }
    return color_map.get(color.lower(), [0, 0, 255])  # Bleu par défaut

def parse_project_prompt(client, prompt, config):
    """
    Analyse un prompt de projet en utilisant un modèle LLM Ollama
    
    Args:
        client: Client Ollama
        prompt (str): Prompt décrivant un projet
        config (dict): Configuration du modèle
    
    Returns:
        dict: Informations structurées du projet
    """
    try:
        # Configuration du modèle
        model = config.get('llm_model', 'mervinpraison/llama3.2-3B-instruct-test-2:8b')
        
        # Prompt système pour l'analyse
        system_prompt = """
        Tu es un assistant spécialisé dans l'analyse de prompts de projet.
        Tu dois identifier si le prompt est une création ou une mise à jour de projet

        Extrait les informations suivantes :
        - Nom du projet
        - Date de début (jour et mois)
        - Date de fin (jour et mois)
        - Couleur du projet
        
        Règles importantes :
        - Le mois de janvier est 0
        - Le mois de décembre est 11
        - Vérifie attentivement le mois de fin

        
        Réponds UNIQUEMENT au format JSON suivant :
        {
            "type": "create|update",
            "task_name": "Nom du projet",
            "start_month": [index_mois, position_dans_mois],
            "start_date": "YYYY/MM/DD",
            "end_month": [index_mois, position_dans_mois],
            "end_date": "YYYY/MM/DD",
            "color_rgb": [R, G, B]
        }
        
        Règles pour la position dans le mois :
        - 0.0 : début du mois (1-10)
        - 0.5 : milieu du mois (11-20)
        - 1.0 : fin du mois (21-31)
        
        Exemples :
        Prompt: "je veux un projet 'P1' du 15 mai au 29 decembre (couleur : vert)"
        Réponse: {
            "type": "create",
            "task_name": "P1",
            "start_month": [4, 0.5],
            "start_date": "2025/05/15",
            "end_month": [11, 1.0],
            "end_date": "2025/12/29",
            "color_rgb": [0, 255, 0]
        }

        Prompt: "j
        je veux un projet 'P1' du 25 mars au 3 aout (couleur : vert)"
        Réponse: {
            "type": "create",
            "task_name": "P1",
            "start_month": [2, 1.0],
            "start_date": "2025/03/25",
            "end_month": [7, 0.0],
            "end_date": "2025/08/03",
            "color_rgb": [0, 255, 0]
        }        

        Prompt: "je veux modifier le projet 'P1' pour qu'il commence le 1er juin"
        Réponse: {
            "type": "update",
            "task_name": "P1",
            "start_month": [5, 0.0],
            "start_date": "2025/06/01",
            "end_month": null,
            "end_date": null,
            "color_rgb": null
        }
        """
        
        # Appel au modèle Ollama
        response = client.chat(
            model=model,
            messages=[
                {'role': 'system', 'content': system_prompt},
                {'role': 'user', 'content': prompt}
            ]
        )
        
        # Extraction du contenu de la réponse
        result = response['message']['content'].strip()

        # Tenter de parser le résultat comme un JSON
        try:
            # Nettoyer le résultat des éventuels caractères superflus
            # Rechercher le premier et le dernier caractère JSON
            start_index = result.find('{')
            end_index = result.rfind('}') + 1
            
            if start_index != -1 and end_index != -1:
                result = result[start_index:end_index]
            
            # Parser le JSON
            parsed_res = json.loads(result)
            
            # Vérifier que le JSON contient les clés requises
            required_keys = ['type', 'task_name']
            if not all(key in parsed_res for key in required_keys):
                raise ValueError("JSON incomplet")
            
            # Gérer les cas de mise à jour et de création
            if parsed_res['type'] == 'update':
                # Supprimer les clés avec des valeurs null
                parsed_res = {k: v for k, v in parsed_res.items() if v is not None}
                                        
            return parsed_res
        
        except json.JSONDecodeError as e:
            print(f"Erreur de décodage JSON : {e}")
            print(f"Contenu problématique : {result}")
            return None
        except ValueError as e:
            print(f"Erreur de validation JSON : {e}")
            return None
    
    except Exception as e:
        print(f"Erreur lors de l'analyse du prompt : {e}")
        return None

def load_config():
    with open('config.yaml', 'r') as f:
        return yaml.safe_load(f)

def encode_image_to_base64(image_path):
    if not os.path.exists(image_path):
        raise FileNotFoundError(f"L'image {image_path} n'existe pas. Veuillez vérifier le chemin.")
    
    # Vérifier la taille et le type de fichier
    file_size = os.path.getsize(image_path)
    file_type = imghdr.what(image_path)
    
    print(f"Fichier image : {image_path}")
    print(f"Taille du fichier : {file_size} octets")
    print(f"Type d'image : {file_type}")
    
    if file_size > 20 * 1024 * 1024:  # Limite à 20 Mo
        raise ValueError("La taille de l'image est trop grande (> 20 Mo)")
    
    if file_type not in ['jpeg', 'png', 'gif', 'bmp']:
        raise ValueError(f"Type de fichier non supporté : {file_type}")
    
    with open(image_path, "rb") as image_file:
        return base64.b64encode(image_file.read()).decode('utf-8')

def analyze_image_with_vision(client, image_path):
    try:
        base64_image = encode_image_to_base64(image_path)
        
        response = client.chat(
            model=config['llm_model'],  
            messages=[
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": "Analyse cette roadmap et donne-moi les informations suivantes en format JSON :\n"
                                   "1. Les mois présents\n"
                                   "2. Les tâches avec leurs périodes\n"
                                   "3. Les couleurs utilisées (en RGB)\n"
                                   "4. La légende (Done, Coming, Not planned)\n"
                        },
                        {
                            "type": "image",
                            "image": base64_image
                        }
                    ]
                }
            ],
            max_tokens=1000
        )
        return response.choices[0].message.content
    except Exception as e:
        print(f"Erreur détaillée lors de l'analyse de l'image : {type(e).__name__} - {str(e)}")
        return None





def create_task_on_roadmap(prs, task_info):
    """
    Crée une tâche sur la roadmap
    
    Args:
        prs (Presentation): Présentation PowerPoint
        task_info (dict): Informations de la tâche
    """
    # Débogage : afficher le contenu complet de task_info
    print("DEBUG - task_info complet :", json.dumps(task_info, indent=2))
    
    # Extraction des informations de la tâche
    task_name = task_info['task_name']
    
    # Convertir les tuples en listes si nécessaire
    start_month = task_info['start_month'][0] if isinstance(task_info['start_month'], (list, tuple)) else task_info['start_month']
    start_pos = task_info['start_month'][1] if isinstance(task_info['start_month'], (list, tuple)) else 0.5
    
    end_month = task_info['end_month'][0] if isinstance(task_info['end_month'], (list, tuple)) else task_info['end_month']
    end_pos = task_info['end_month'][1] if isinstance(task_info['end_month'], (list, tuple)) else 1.0
    
    # Valeurs par défaut si None
    start_month = 0 if start_month is None else start_month
    end_month = 11 if end_month is None else end_month
    
    # Valeurs par défaut pour les positions
    start_pos = 0.5 if start_pos is None else start_pos
    end_pos = 1.0 if end_pos is None else end_pos
    
    color_rgb = task_info['color_rgb']         # Couleur RGB

    # Récupération du slide de roadmap
    roadmap_slide = prs.slides[0]  # Première slide (roadmap)

    # Dimensions de la slide
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Dimensions de la grille
    grid_margin_top = Inches(1.5)  # Marge en haut
    grid_margin_bottom = Inches(0.5)  # Marge en bas
    grid_margin_left = Inches(0.5)  # Marge à gauche
    grid_margin_right = Inches(0.5)  # Marge à droite

    # Calculer la hauteur de la grille
    grid_height = slide_height - grid_margin_top - grid_margin_bottom
    grid_width = slide_width - grid_margin_left - grid_margin_right

    # Calculer la largeur d'une colonne de mois
    month_width = grid_width / 12

    # Calculer la position horizontale de début et de fin
    print(f"start_month : {start_month}, start_pos : {start_pos}")
    print(f"end_month : {end_month}, end_pos : {end_pos}")
    start_x = grid_margin_left + (start_month * month_width) + (start_pos * month_width)
    end_x = grid_margin_left + (end_month * month_width) + (end_pos * month_width)

    # Hauteur de la tâche
    task_height = Inches(0.5)
    
    # Calculer la position verticale dynamique
    existing_shapes = [shape for shape in roadmap_slide.shapes if shape.has_text_frame]
    
    # Décaler vers le bas (0.5 inch sous l'axe des mois)
    task_y = grid_margin_top + Inches(0.5) + (len(existing_shapes) * Inches(0.6))

    # Créer la forme de la tâche
    task_shape = roadmap_slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 
        start_x, 
        task_y, 
        end_x - start_x, 
        task_height
    )

    # Formater la forme
    task_shape.fill.solid()
    task_shape.fill.fore_color.rgb = RGBColor(color_rgb[0], color_rgb[1], color_rgb[2])
    task_shape.line.fill.background()

    # Ajouter le texte de la tâche
    text_frame = task_shape.text_frame
    text_frame.text = task_name
    text_frame.paragraphs[0].font.size = Pt(10)
    text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # Texte en noir
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def create_roadmap_slide(prs, task_info=None):
    """Crée ou met à jour un slide de roadmap"""
    print("\n=== Création/Mise à jour du slide de roadmap ===")
    
    # Utiliser le premier slide ou en créer un nouveau
    if len(prs.slides) > 0:
        slide = prs.slides[0]
        print(f"Utilisation du slide existant (index 0)")
    else:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Layout vide
        print(f"Création d'un nouveau slide")
    
    # Vérifier si la grille des mois existe déjà
    months_grid_exists = any(shape.has_table for shape in slide.shapes)
    
    # Titre "ROADMAP"
    if not any(shape.has_text_frame and shape.text_frame.text == "ROADMAP" for shape in slide.shapes):
        title = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(0.5))
        title.text_frame.text = "ROADMAP"
        title.text_frame.paragraphs[0].font.size = Pt(24)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)    
    
    if not months_grid_exists:
        print("Ajout de la grille des mois")
        
        # Calculer les dimensions de la slide
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        
        print(f"Dimensions de la slide : {slide_width / 914400:.2f} x {slide_height / 914400:.2f} inches")
        
        # Marges internes
        margin_left = Inches(0.5)
        margin_right = Inches(0.5)
        
        # Calculer la largeur effective
        effective_width = slide_width - (margin_left + margin_right)
        
        # Créer la grille des mois
        months_box = slide.shapes.add_table(
            2,  # 2 rangées
            12,  # 12 colonnes (mois)
            margin_left,  # Position X de départ
            Inches(1.5),  # Position Y
            effective_width,  # Largeur totale
            Inches(0.5)  # Hauteur
        ).table
        
        months = ["JAN", "FEB", "MAR", "APR", "MAY", "JUN", "JUL", "AUG", "SEP", "OCT", "NOV", "DEC"]
        
        for i, month in enumerate(months):
            cell = months_box.cell(0, i)
            cell.text = month
            cell.text_frame.paragraphs[0].font.size = Pt(8)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    else:
        print("Grille des mois déjà existante")
    
    # Ajouter la nouvelle tâche si spécifiée
    if task_info:
        create_task_on_roadmap(prs, task_info)
    
    return slide

def convert_db_task_to_task_info(db_task):
    """
    Convertit une tâche de la base de données en task_info.
    
    Args:
        db_task (dict): Tâche récupérée de la base de données
    
    Returns:
        dict: Tâche au format task_info
    """
    # Convertir color_rgb de JSON à liste si nécessaire
    color_rgb = json.loads(db_task['color_rgb']) if db_task['color_rgb'] else None
    
    task_info = {
        "type": "create",  # Par défaut, toujours "create"
        "task_name": db_task['task_name'],
        "start_month": [
            db_task['start_month'] if db_task['start_month'] is not None else None,
            db_task['start_position'] if db_task['start_position'] is not None else 0.5
        ],
        "end_month": [
            db_task['end_month'] if db_task['end_month'] is not None else None,
            db_task['end_position'] if db_task['end_position'] is not None else 1.0
        ],
        "color_rgb": color_rgb
    }
    
    return task_info

def normalize_text(text):
    """
    Normalise un texte en le nettoyant et le standardisant.
    
    Args:
        text (str): Texte à normaliser
    
    Returns:
        str: Texte normalisé
    """
    # Convertir en minuscules
    text = text.lower()
    
    # Supprimer les espaces en début et fin
    text = text.strip()
    
    # Remplacer les caractères spéciaux et multiples espaces
    import re
    text = re.sub(r'\s+', ' ', text)  # Remplacer les espaces multiples par un seul
    text = re.sub(r'[^\w\s]', '', text)  # Supprimer la ponctuation
    
    return text

def list_powerpoint_objects(file_path):
    """
    Liste tous les objets d'un fichier PowerPoint avec leurs labels.
    
    Args:
        file_path (str): Chemin complet vers le fichier PowerPoint
    
    Returns:
        list: Liste des textes normalisés trouvés dans le fichier PowerPoint
    """
    # Vérifier si le fichier existe
    if not os.path.exists(file_path):
        print(f"Erreur : Le fichier {file_path} n'existe pas.")
        return []

    # Charger la présentation
    prs = Presentation(file_path)
    
    # Liste des objets à ignorer
    ignored_objects = ['ROADMAP']
    
    # Liste pour stocker les objets
    objects_list = []
    
    # Parcourir tous les slides
    for slide_index, slide in enumerate(prs.slides, 1):
        slide_objects = []
        
        # Parcourir les formes de chaque slide
        for shape_index, shape in enumerate(slide.shapes, 1):
            # Vérifier si la forme a un cadre de texte
            if shape.has_text_frame:
                # Extraire le texte du cadre
                text = shape.text_frame.text.strip()
                
                # Ajouter le label si non vide et non ignoré
                if text and text not in ignored_objects:
                    slide_objects.append({
                        'type': 'texte',
                        'slide': slide_index,
                        'index': shape_index,
                        'text': text
                    })
            
            # Vérifier si c'est un tableau
            elif shape.has_table:
                # Ignorer le tableau des mois
                table = shape.table
                first_row_texts = [cell.text.strip() for cell in table.rows[0].cells]
                if not (len(first_row_texts) > 1 and all(len(month) == 3 for month in first_row_texts)):
                    table_rows = []
                    for row_index, row in enumerate(table.rows, 1):
                        row_texts = [cell.text.strip() for cell in row.cells]
                        if any(row_texts):
                            table_rows.append({
                                'row_index': row_index,
                                'row_texts': row_texts
                            })
                    
                    slide_objects.append({
                        'type': 'tableau',
                        'slide': slide_index,
                        'index': shape_index,
                        'rows': table_rows
                    })
        
        # Ajouter les objets du slide à la liste principale si non vide
        if slide_objects:
            objects_list.extend(slide_objects)
    
    # Normaliser et retourner uniquement les textes
    return [normalize_text(obj['text']) for obj in objects_list if obj['type'] == 'texte']

def main():
    # Créer les dossiers de sortie
    templates_dir = "templates"
    output_dir = "generated"
    
    # Créer les dossiers s'ils n'existent pas
    os.makedirs(templates_dir, exist_ok=True)
    os.makedirs(output_dir, exist_ok=True)
    
    # Chemins complets
    template_path = os.path.join(templates_dir, "roadmap.pptx")
    output_path = os.path.join(output_dir, "roadmap.pptx")
    
    # Supprimer le fichier existant dans le répertoire de génération
    if os.path.exists(output_path):
        try:
            os.remove(output_path)
            print(f"Fichier existant supprimé : {output_path}")
        except Exception as e:
            print(f"Erreur lors de la suppression du fichier : {e}")
    
    # Charger ou créer la présentation
    if os.path.exists(template_path):
        prs = Presentation(template_path)
        print(f"Chargement du template depuis {template_path}")
    else:
        prs = Presentation()
        print("Création d'une nouvelle présentation")
        
        # Sauvegarde du template vide
        prs.save(template_path)
        print(f"Template vide sauvegardé dans {template_path}")
    
    # Configuration
    config = load_config()
    client = ollama.Client(host=config.get('ollama_host', 'http://localhost:11434'))
    
    # Traitement des prompts
    prompts = []
    
    # Vérifier si des arguments en ligne de commande sont passés
    if len(sys.argv) > 1:
        # Concaténer tous les arguments en ligne de commande
        prompts = [' '.join(sys.argv[1:])]
    
    # Si pas d'arguments en ligne de commande, essayer de lire le fichier
    if not prompts:
        try:
            with open('prompt.txt', 'r') as f:
                prompts = [line.strip() for line in f.readlines() if line.strip()]
        except FileNotFoundError:
            # Si le fichier n'existe pas, demander une saisie interactive
            prompts = [input("Entrez votre demande (ex: ajoute le projet TOTO se déroule du 1er février au 13 juin (couleur : orange)) : ")]
    
    # Traiter chaque prompt successivement
    for prompt in prompts:
        print(f"\n--- Traitement du prompt : {prompt} ---")
        
        try:
            # Analyser le prompt
            task_info = parse_project_prompt(client, prompt, config)
            
            if task_info:
                # Insérer ou mettre à jour la tâche dans la base de données
                task_id = task_db.upsert_task(task_info, raw_prompt=prompt)
                print(f"Tâche créée ou mise à jour avec l'ID : {task_id}")
                
            else:
                print("Impossible de parser le prompt")
        
        except Exception as e:
            print(f"Erreur lors du traitement du prompt '{prompt}' : {e}")
            traceback.print_exc()
    
    # Récupérer toutes les tâches de la base de données
    all_tasks = task_db.list_tasks()
    
    # Réinitialiser la présentation
    prs = Presentation(template_path)
    
    # Recréer la slide avec toutes les tâches
    for task in all_tasks:
        # Convertir la tâche de la base de données en task_info
        task_info = convert_db_task_to_task_info(task)
        
        # Créer ou mettre à jour le slide de roadmap
        create_roadmap_slide(prs, task_info)
    
    # Sauvegarder la présentation finale
    prs.save(output_path)
    print(f"Présentation finale mise à jour : {output_path}")
    
    print("\nTraitement de tous les prompts terminé.")

if __name__ == "__main__":
    main()