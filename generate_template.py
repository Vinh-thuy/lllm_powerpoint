from pptx import Presentation

prs = Presentation()

# Slide titre
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "{{title}}"
slide.placeholders[1].text = "{{subtitle}}"

# Slide planning
slide = prs.slides.add_slide(prs.slide_layouts[1])
table = slide.shapes.add_table(rows=3, cols=3, left=100, top=200, width=500, height=300).table
table.cell(0, 0).text = "{{schedule_header}}"

# Slide contenu
slide = prs.slides.add_slide(prs.slide_layouts[2])
slide.shapes.add_textbox(100, 100, 300, 150).text = "{{content}}"

prs.save('templates/base_template.pptx')
