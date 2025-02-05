import os
import collections.abc
from pptx import Presentation

def list_powerpoint_objects(file_path):
    """
    Liste tous les objets d'un fichier PowerPoint avec leurs labels.
    
    Args:
        file_path (str): Chemin complet vers le fichier PowerPoint
    """
    # Vérifier si le fichier existe
    if not os.path.exists(file_path):
        print(f"Erreur : Le fichier {file_path} n'existe pas.")
        return

    # Charger la présentation
    prs = Presentation(file_path)
    
    # Liste des objets à ignorer
    ignored_objects = ['ROADMAP']
    
    # Parcourir tous les slides
    for slide_index, slide in enumerate(prs.slides, 1):
        print(f"\n--- Slide {slide_index} ---")
        
        # Parcourir les formes de chaque slide
        for shape_index, shape in enumerate(slide.shapes, 1):
            # Vérifier si la forme a un cadre de texte
            if shape.has_text_frame:
                # Extraire le texte du cadre
                text = shape.text_frame.text.strip()
                
                # Afficher le label si non vide et non ignoré
                if text and text not in ignored_objects:
                    print(f"  Objet {shape_index}: {text}")
            
            # Vérifier si c'est un tableau
            elif shape.has_table:
                # Ignorer le tableau des mois
                table = shape.table
                first_row_texts = [cell.text.strip() for cell in table.rows[0].cells]
                if not (len(first_row_texts) > 1 and all(len(month) == 3 for month in first_row_texts)):
                    print(f"  Tableau {shape_index}")
                    for row_index, row in enumerate(table.rows, 1):
                        row_texts = [cell.text.strip() for cell in row.cells]
                        if any(row_texts):
                            print(f"    Ligne {row_index}: {' | '.join(row_texts)}")

def main():
    # Chemin du fichier PowerPoint à analyser
    powerpoint_path = "/Users/vinh/Documents/LLM_PPT/generated/roadmap.pptx"
    
    # Lister les objets
    list_powerpoint_objects(powerpoint_path)

if __name__ == "__main__":
    main()
