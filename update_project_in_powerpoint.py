import json
import re
from pptx import Presentation

def get_precise_date(month_info):
    """
    Convertit les informations de mois en date précise.
    
    Args:
        month_info (list): [index_mois, position_dans_mois]
    
    Returns:
        str: Date précise
    """
    mois = [
        'janvier', 'février', 'mars', 'avril', 'mai', 'juin', 
        'juillet', 'août', 'septembre', 'octobre', 'novembre', 'décembre'
    ]
    
    if not isinstance(month_info, list) or len(month_info) != 2:
        return None
    
    index_mois, position = month_info
    
    if index_mois < 0 or index_mois > 11:
        return None
    
    # Déterminer le jour précis
    if position == 0.0:
        jour = '1'
    elif position == 0.5:
        jour = '15'
    elif position == 1.0:
        # Dernier jour du mois
        jours_par_mois = [31, 28, 31, 30, 31, 30, 31, 31, 30, 31, 30, 31]
        jour = str(jours_par_mois[index_mois])
    else:
        return None
    
    return f"{jour} {mois[index_mois]}"

def update_project_in_powerpoint(file_path, update_info):
    """
    Met à jour un projet dans un fichier PowerPoint.
    
    Args:
        file_path (str): Chemin du fichier PowerPoint
        update_info (dict): Informations de mise à jour du projet
    
    Returns:
        bool: True si la mise à jour a réussi, False sinon
    """
    # Vérifier les informations de mise à jour
    if not update_info or 'type' not in update_info or update_info['type'] != 'update':
        print("Erreur : Informations de mise à jour invalides")
        return False
    
    # Charger la présentation
    prs = Presentation(file_path)
    
    # Parcourir tous les slides
    for slide in prs.slides:
        # Parcourir toutes les formes du slide
        for shape in slide.shapes:
            # Vérifier si la forme a un cadre de texte
            if shape.has_text_frame:
                # Extraire le texte du cadre
                text = shape.text_frame.text.strip()
                
                # Vérifier si le texte correspond au nom du projet
                if text.startswith(update_info.get('task_name')):
                    # Extraire les dates actuelles
                    date_match = re.search(r'(\d{1,2}\s+\w+)\s*-\s*(\d{1,2}\s+\w+)', text)
                    
                    if date_match:
                        current_start_date = date_match.group(1)
                        current_end_date = date_match.group(2)
                        
                        # Déterminer les nouvelles dates
                        new_start_date = (get_precise_date(update_info.get('start_month')) 
                                          if 'start_month' in update_info 
                                          else current_start_date)
                        new_end_date = (get_precise_date(update_info.get('end_month')) 
                                        if 'end_month' in update_info 
                                        else current_end_date)
                        
                        # Remplacer le texte avec les nouvelles dates
                        new_text = text.replace(
                            f"{current_start_date} - {current_end_date}", 
                            f"{new_start_date} - {new_end_date}"
                        )
                        
                        # Mettre à jour le texte
                        shape.text_frame.text = new_text
                        
                        print(f"Projet '{update_info['task_name']}' mis à jour : {new_start_date} - {new_end_date}")
                        
                        # Sauvegarder la présentation
                        prs.save(file_path)
                        return True
    
    print(f"Projet '{update_info.get('task_name')}' non trouvé")
    return False

def main():
    # Chemin du fichier PowerPoint
    file_path = '/Users/vinh/Documents/LLM_PPT/templates/roadmap.pptx'
    
    # Exemples de mise à jour
    update_info_1 = {
        "type": "update",
        "task_name": "P1",
        "start_month": [2, 0.5]  
    }
    
    # Tester différents scénarios de mise à jour
    update_project_in_powerpoint(file_path, update_info_1)
   

if __name__ == "__main__":
    main()
